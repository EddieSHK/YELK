import os
# Set the environment variable to use the pure-Python implementation of protobuf
os.environ['PROTOCOL_BUFFERS_PYTHON_IMPLEMENTATION'] = 'python'
# Set the environment variable to avoid OpenMP errors
os.environ["KMP_DUPLICATE_LIB_OK"] = "TRUE"
from flask import Flask, render_template, request, jsonify
from PyPDF2 import PdfReader
import google.generativeai as genai
import json
from PIL import Image
from io import BytesIO
from docx import Document
from pptx import Presentation
from paddleocr import PaddleOCR

app = Flask(__name__)

# Track questions and flashcards
TRACKED_FILE = "tracked_questions.json"


def load_tracked():
    if os.path.exists(TRACKED_FILE):
        with open(TRACKED_FILE, 'r') as file:
            return json.load(file)
    return {"questions": [], "flashcards": []}


tracked = load_tracked()


def save_tracked():
    with open(TRACKED_FILE, 'w') as file:
        json.dump(tracked, file)


def extract_text_from_pdf(file):
    """Extract text from an uploaded PDF file."""
    reader = PdfReader(file)
    return " ".join([page.extract_text() for page in reader.pages])


def extract_text_from_docx(file):
    """Extract text from a Word document."""
    doc = Document(file)
    return " ".join([para.text for para in doc.paragraphs])


def extract_text_from_pptx(file):
    """Extract text from a PowerPoint presentation."""
    presentation = Presentation(file)
    text = []
    for slide in presentation.slides:
        for shape in slide.shapes:
            if hasattr(shape, 'text'):
                text.append(shape.text)
    return " ".join(text)


def extract_text_from_image(file):
    """Extract text from an image using PaddleOCR."""
    try:
        # Initialize PaddleOCR (ensure that paddleocr is installed and configured correctly)
        ocr = PaddleOCR(use_angle_cls=True, lang='en')  # Set language as English

        # Ensure file is opened correctly
        image = Image.open(file)

        # Convert the image to RGB if it's not already (PaddleOCR requires RGB)
        if image.mode != 'RGB':
            image = image.convert('RGB')

        # Perform OCR on the image
        result = ocr.ocr(image, cls=True)

        # Extract text from the OCR result
        text = ""
        for line in result[0]:
            text += line[1] + " "
        return text
    except Exception as e:
        print(f"Error processing image: {e}")
        return ""


def ask_gemini_ai_with_genai(text: str):
    """Generate a summary and flashcards using Google Generative AI."""
    genai.configure(api_key="AIzaSyADe46GEGmju9hg8KMIc3qj2LuQ21ohiDE")
    model = genai.GenerativeModel('gemini-pro')

    # Generate content
    prompt = (
        f"Summarize the following content in paragraph form :\n\n{text}"
    )
    prompt2 = (
        f"Generate flashcard with question and answer: in paragraph form\n\n{text}"
    )
    response = model.generate_content(prompt)
    response2 = model.generate_content(prompt2)

    # Parse the generated content
    generated_text = response.text
    print(generated_text)
    summary = [response.text]
    flashcards = [response2.text]

    print("Summary:", summary)
    print("FLASHCARD:", flashcards)
    return summary, flashcards


def ask_gemini_for_answer(question: str):
    """Query Gemini AI for an answer to a specific question."""
    genai.configure(api_key="AIzaSyADe46GEGmju9hg8KMIc3qj2LuQ21ohiDE")
    model = genai.GenerativeModel('gemini-pro')

    # Generate answer to the question
    prompt = f"Answer the following question:\n\n{question}"

    response = model.generate_content(prompt)

    # Return the response text
    return response.text

@app.route('/')
def index():
    return render_template('index2.html')


@app.route('/ask_question', methods=['POST'])
def ask_question():
    """Endpoint for asking a custom question to Gemini AI."""
    question = request.form.get('question')

    if not question:
        return jsonify({"error": "No question provided"}), 400

    # Get the answer from Gemini AI
    answer = ask_gemini_for_answer(question)

    return jsonify({
        "question": question,
        "answer": answer
    })
@app.route('/generate_flashcards', methods=['POST'])

def generate_flashcards():
    if 'files[]' not in request.files:
        return jsonify({"error": "No files uploaded"}), 400

    files = request.files.getlist('files[]')  # Get all uploaded files
    if not files or all(file.filename == '' for file in files):
        return jsonify({"error": "No selected files"}), 400

    summaries = []
    flashcards = []

    for file in files:
        if file.filename == '':
            continue

        # Process each file based on its type
        if file.filename.endswith('.pdf'):
            text = extract_text_from_pdf(file)
        elif file.filename.endswith('.docx'):
            text = extract_text_from_docx(file)
        elif file.filename.endswith('.pptx'):
            text = extract_text_from_pptx(file)
        elif file.filename.lower().endswith(('.png', '.jpg', '.jpeg')):
            text = extract_text_from_image(file)
        else:
            continue  # Skip unsupported file formats

        # Process text in chunks if it exceeds token limits
        file_summary = []
        file_flashcards = []
        for chunk in split_text_into_chunks(text, 2000): 
            summary, chunk_flashcards = ask_gemini_ai_with_genai(chunk)
            file_summary.append(summary)
            file_flashcards.extend(chunk_flashcards)

        # # Generate summary and flashcards for the current file
        # file_summary, file_flashcards = ask_gemini_ai_with_genai(text)

        # Append the processed file results
        summaries.append(file_summary)
        flashcards.append(file_flashcards)

    # # Flatten the lists of summaries and flashcards to combine results
    # combined_summary = [item for sublist in summaries for item in sublist]
    # combined_flashcards = [item for sublist in flashcards for item in sublist]

    # Flatten the lists of summaries and flashcards to combine results
    combined_summary = [truncate_text(item,2000) for sublist in summaries for item in sublist]
    combined_flashcards = [truncate_text(item,2000) for sublist in flashcards for item in sublist]

    # Save to tracked data (local storage or database)
    tracked["summary"] = tracked.get("summary", []) + combined_summary
    tracked["flashcards"] = tracked.get("flashcards", []) + combined_flashcards
    save_tracked()

    return jsonify({
        "summary_points": combined_summary,
        "flashcards": combined_flashcards,
        "store_in_local": True  # Suggests storing data in local storage for subsequent use
    })

def split_text_into_chunks(text, max_token_limit=4000):
    # raise RuntimeError("split_text_into_chunks should not be called!")
    words = text.split()
    for i in range(0, len(words), max_token_limit):
        yield " ".join(words[i:i + max_token_limit])

def truncate_text(text, max_tokens=4000):
    words = text.split()
    return " ".join(words[:max_tokens]) + ("..." if len(words) > max_tokens else "")

if __name__ == "__main__":
    app.run(debug=True)
