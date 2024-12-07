import os
# Set the environment variable to use the pure-Python implementation of protobuf
os.environ['PROTOCOL_BUFFERS_PYTHON_IMPLEMENTATION'] = 'python'
# Set the environment variable to avoid OpenMP errors
os.environ["KMP_DUPLICATE_LIB_OK"] = "TRUE"
from flask import Flask, render_template, request, jsonify, send_file
from PyPDF2 import PdfReader
import google.generativeai as genai
import json
from PIL import Image
from io import BytesIO
from docx import Document
from pptx import Presentation
from paddleocr import PaddleOCR
import numpy as np
from fpdf import FPDF


app = Flask(__name__)

# Track questions and flashcards
TRACKED_FILE = "tracked_questions.json"


def load_tracked():
    if os.path.exists(TRACKED_FILE):
        with open(TRACKED_FILE, 'r') as file:
            return json.load(file)
    return {"questions": [], "flashcards": []}


tracked = load_tracked()


def save_tracked():
    with open(TRACKED_FILE, 'w') as file:
        json.dump(tracked, file)


def extract_text_from_pdf(file):
    """Extract text from an uploaded PDF file."""
    reader = PdfReader(file)
    return " ".join([page.extract_text() for page in reader.pages])


def extract_text_from_docx(file):
    """Extract text from a Word document."""
    doc = Document(file)
    return " ".join([para.text for para in doc.paragraphs])


def extract_text_from_pptx(file):
    """Extract text from a PowerPoint presentation."""
    presentation = Presentation(file)
    text = []
    for slide in presentation.slides:
        for shape in slide.shapes:
            if hasattr(shape, 'text'):
                text.append(shape.text)
    return " ".join(text)


import easyocr

def extract_text_from_image(file):
    """Extract text from an image using EasyOCR."""
    try:
        # Read the file content into bytes
        file_bytes = file.read()  # Read the uploaded file into memory

        # Convert the bytes into a PIL image
        image = Image.open(BytesIO(file_bytes))

        # Convert the PIL Image to a numpy array (required by EasyOCR)
        image_np = np.array(image)

        # Initialize EasyOCR reader (English language)
        reader = easyocr.Reader(['en'])

        # Use EasyOCR to read the text from the image
        result = reader.readtext(image_np)

        # Extract and join the text from the result
        text = " ".join([item[1] for item in result])

        if not text.strip():
            print("No text detected in the image.")
        
        return text.strip()

    except Exception as e:
        print(f"Error processing image: {e}")
        return ""

def ask_gemini_ai_with_genai(text: str):
    """Generate a summary and flashcards using Google Generative AI."""
    genai.configure(api_key="AIzaSyADe46GEGmju9hg8KMIc3qj2LuQ21ohiDE")
    model = genai.GenerativeModel('gemini-pro')

    # Generate content
    prompt = (
        f"Summarize the following content in paragraph form :\n\n{text}"
    )
    prompt2 = (
        f"Generate flashcard with question and answer: in paragraph form\n\n{text}"
    )
    response = model.generate_content(prompt)
    response2 = model.generate_content(prompt2)

    # Parse the generated content
    generated_text = response.text
    print(generated_text)
    summary = [response.text]
    flashcards = [response2.text]

    print("Summary:", summary)
    print("FLASHCARD:", flashcards)
    return summary, flashcards

def ask_gemini_ai_with_genai_summary(text: str):
    """Generate a summary using Google Generative AI."""
    genai.configure(api_key="AIzaSyADe46GEGmju9hg8KMIc3qj2LuQ21ohiDE")
    model = genai.GenerativeModel('gemini-pro')

    # Generate content
    prompt = (
        f"Summarize the following content in paragraph form :\n\n{text}"
    )

    response = model.generate_content(prompt)

    # Parse the generated content
    generated_text = response.text
    print(generated_text)
    summary = [response.text]

    print("Summary:", summary)
    return summary

def ask_gemini_ai_with_genai_flashCard(text: str):
    """Generate a flashcards using Google Generative AI."""
    genai.configure(api_key="AIzaSyADe46GEGmju9hg8KMIc3qj2LuQ21ohiDE")
    model = genai.GenerativeModel('gemini-pro')

    # Generate content
    prompt2 = (
        f"Generate 5 flashcard question using the text i want just the question without any numbering or symbol:\n\n{text}"
    )

    response2 = model.generate_content(prompt2)
    questions = response2.text.strip().split("\n")  
    #flashcards = [response2.text]
    flashcards = []

    for question in questions:
        if question.strip():  # Ensure the question is not empty
            # Step 2: Generate answers for each question
            answer_prompt = f"Answer the following question based on the text without extra thing or symbol or numbering just the answer:\n\nText:\n{text}\n\nQuestion:\n{question}"
            answer_response = model.generate_content(answer_prompt)

            # Construct the flashcard
            flashcards.append({
                "question": question.strip(),
                "answer": answer_response.text.strip()
            })
    
    print("FLASHCARD:", flashcards)
    return flashcards

def ask_gemini_for_answer(question: str):
    """Query Gemini AI for an answer to a specific question."""
    genai.configure(api_key="AIzaSyADe46GEGmju9hg8KMIc3qj2LuQ21ohiDE")
    model = genai.GenerativeModel('gemini-pro')

    # Generate answer to the question
    prompt = f"Answer the following question:\n\n{question}"

    response = model.generate_content(prompt)

    # Return the response text
    return response.text

@app.route('/')
def index():
    return render_template('index2.html')


@app.route('/ask_question', methods=['POST'])
def ask_question():
    """Endpoint for asking a custom question to Gemini AI."""
    question = request.form.get('question')

    if not question:
        return jsonify({"error": "No question provided"}), 400

    # Get the answer from Gemini AI
    answer = ask_gemini_for_answer(question)

    return jsonify({
        "question": question,
        "answer": answer
    })

@app.route('/generate_flashcards', methods=['POST'])
def generate_flashcards():
    """Endpoint to generate flashcards based on user-selected number of questions."""
    question_count = int(request.form.get('question_count', 3))  # Default to 3 questions

    if 'files[]' not in request.files:
        return jsonify({"error": "No files uploaded"}), 400

    files = request.files.getlist('files[]')  # Get all uploaded files
    if not files or all(file.filename == '' for file in files):
        return jsonify({"error": "No selected files"}), 400

    flashcards = []

    for file in files:
        if file.filename == '':
            continue

        # Process each file based on its type
        if file.filename.endswith('.pdf'):
            text = extract_text_from_pdf(file)
        elif file.filename.endswith('.docx'):
            text = extract_text_from_docx(file)
        elif file.filename.endswith('.pptx'):
            text = extract_text_from_pptx(file)
        elif file.filename.lower().endswith(('.png', '.jpg', '.jpeg')):
            text = extract_text_from_image(file)
        else:
            continue  # Skip unsupported file formats

        # Process text in chunks if it exceeds token limits
        for chunk in split_text_into_chunks(text, 2000):
            chunk_flashcards = ask_gemini_ai_with_genai_flashCard(chunk)  # Only get flashcards
            flashcards.extend(chunk_flashcards)

    # Limit the flashcards to the user-selected number
    limited_flashcards = flashcards[:question_count]

    # Save to tracked data
    tracked["flashcards"] = tracked.get("flashcards", []) + limited_flashcards
    save_tracked()

    return jsonify({
        "flashcards": limited_flashcards,
        "store_in_local": True  # Suggest storing data in local storage
    })




@app.route('/generate_summary', methods=['POST'])
def generate_summary():
    """Endpoint to generate summaries."""
    if 'files[]' not in request.files:
        return jsonify({"error": "No files uploaded"}), 400

    files = request.files.getlist('files[]')  # Get all uploaded files
    if not files or all(file.filename == '' for file in files):
        return jsonify({"error": "No selected files"}), 400

    summaries = []

    for file in files:
        if file.filename == '':
            continue

        # Process each file based on its type
        if file.filename.endswith('.pdf'):
            text = extract_text_from_pdf(file)
        elif file.filename.endswith('.docx'):
            text = extract_text_from_docx(file)
        elif file.filename.endswith('.pptx'):
            text = extract_text_from_pptx(file)
        elif file.filename.lower().endswith(('.png', '.jpg', '.jpeg')):
            text = extract_text_from_image(file)
        else:
            continue  # Skip unsupported file formats

        # Process text in chunks if it exceeds token limits
        file_summary = []
        for chunk in split_text_into_chunks(text, 2000):
            summary= ask_gemini_ai_with_genai_summary(chunk)  # Only get the summary
            file_summary.append(summary)

        # Append the processed file results
        summaries.append(file_summary)

    # Flatten the list of summaries
    combined_summary = [
        truncate_text(item, 2000)
        for sublist1 in summaries
        for sublist2 in sublist1
        for item in sublist2
        if isinstance(item, str)  # Ensure only strings are processed
    ]

    # Save to tracked data (local storage or database)
    tracked["summary"] = tracked.get("summary", []) + combined_summary
    save_tracked()

    return jsonify({
        "summary_points": combined_summary,
        "store_in_local": True  # Suggests storing data in local storage for subsequent use
    })

def split_text_into_chunks(text, max_token_limit=4000):
    # raise RuntimeError("split_text_into_chunks should not be called!")
    words = text.split()
    for i in range(0, len(words), max_token_limit):
        yield " ".join(words[i:i + max_token_limit])

def truncate_text(text, max_tokens=4000):
    words = text.split()
    return " ".join(words[:max_tokens]) + ("..." if len(words) > max_tokens else "")


# @app.route('/download_summary_pptx', methods=['POST'])
# def handle_summary_pptx():
#     # Receive and process data for PPTX
#     data = request.get_json()
#     if not data or 'summary_points' not in data:
#         return jsonify({"error": "Invalid data"}), 400

#     # Extract summary points
#     summary_points = data['summary_points']
#     print("Received summary for PPTX:", summary_points)

#     # Respond with success message
#     return jsonify({"message": "Summary data received for PPTX"}), 200
# Create summary PDF function

def create_pdf(title, summary_points):
    pdf = FPDF()
    pdf.set_auto_page_break(auto=True, margin=15)
    pdf.add_page()
    pdf.set_font("Arial", size=12)

    # Add the title
    pdf.set_font("Arial", style="B", size=16)
    pdf.cell(200, 10, txt=title, ln=True, align="C")
    pdf.ln(10)  # Line break

    pdf.set_font("Arial", size=12)
    # Add the summary points
    max_width = 190  # Width of the page minus margins

    # Combine all the points into one block of text
    combined_text = "\n\n".join(summary_points)

    # Write the combined text into the PDF
    pdf.multi_cell(max_width, 10, combined_text)

    # Use BytesIO to save the PDF into memory
    pdf_buffer = BytesIO()
    pdf_content = pdf.output(dest='S').encode('latin1')  # Get the PDF content as bytes
    pdf_buffer.write(pdf_content)
    pdf_buffer.seek(0)  # Rewind the buffer to the beginning
    return pdf_buffer

@app.route("/download_summary_pdf", methods=["POST"])
def handle_summary_pdf():
    summary_points = request.json.get("summary_points", [])
    pdf_file = create_pdf("Slide Summary", summary_points)
    
    return send_file(
        pdf_file,
        as_attachment=True,
        download_name="summary.pdf",
        mimetype="application/pdf"
    )

# @app.route('/download_flashcard_pptx', methods=['POST'])
# def handle_flashcard_pptx():
#     # Receive and process data for PPTX
#     data = request.get_json()
#     if not data or 'flashcard_points' not in data:
#         return jsonify({"error": "Invalid data"}), 400

#     # Extract flashcard points
#     flashcard_points = data['flashcard_points']
#     print("Received flashcard for PPTX:", flashcard_points)

#     # Respond with success message
#     return jsonify({"message": "Flashcard data received for PPTX"}), 200

@app.route('/download_flashcard_pdf', methods=['POST'])
def handle_flashcard_pdf():
    # Receive and process data for PPTX
    data = request.get_json()
    if not data or 'flashcard_points' not in data:
        return jsonify({"error": "Invalid data"}), 400

    # Extract flashcard points
    flashcard_points = data['flashcard_points']
    print("Received flashcard for PDF:", flashcard_points)

    # Generate the PDF
    pdf_file = create_pdf("Flashcards",flashcard_points)

    # Serve the generated PDF as a downloadable file
    return send_file(
        pdf_file,
        as_attachment=True,
        download_name="flashcard.pdf",
        mimetype='application/pdf')

if __name__ == "__main__":
    app.run(debug=True)
