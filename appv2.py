import os
# Set the environment variable to use the pure-Python implementation of protobuf
os.environ['PROTOCOL_BUFFERS_PYTHON_IMPLEMENTATION'] = 'python'
# Set the environment variable to avoid OpenMP errors
os.environ["KMP_DUPLICATE_LIB_OK"] = "TRUE"
from flask import Flask, render_template, request, jsonify
from PyPDF2 import PdfReader
import google.generativeai as genai
import json
from PIL import Image
from io import BytesIO
from docx import Document
from pptx import Presentation
from paddleocr import PaddleOCR

app = Flask(__name__)

# Track questions and flashcards
TRACKED_FILE = "tracked_questions.json"


def load_tracked():
    if os.path.exists(TRACKED_FILE):
        with open(TRACKED_FILE, 'r') as file:
            return json.load(file)
    return {"questions": [], "flashcards": []}


tracked = load_tracked()


def save_tracked():
    with open(TRACKED_FILE, 'w') as file:
        json.dump(tracked, file)


def extract_text_from_pdf(file):
    """Extract text from an uploaded PDF file."""
    reader = PdfReader(file)
    return " ".join([page.extract_text() for page in reader.pages])


def extract_text_from_docx(file):
    """Extract text from a Word document."""
    doc = Document(file)
    return " ".join([para.text for para in doc.paragraphs])


def extract_text_from_pptx(file):
    """Extract text from a PowerPoint presentation."""
    presentation = Presentation(file)
    text = []
    for slide in presentation.slides:
        for shape in slide.shapes:
            if hasattr(shape, 'text'):
                text.append(shape.text)
    return " ".join(text)


def extract_text_from_image(file):
    """Extract text from an image using PaddleOCR."""
    try:
        # Initialize PaddleOCR (ensure that paddleocr is installed and configured correctly)
        ocr = PaddleOCR(use_angle_cls=True, lang='en')  # Set language as English

        # Ensure file is opened correctly
        image = Image.open(file)

        # Convert the image to RGB if it's not already (PaddleOCR requires RGB)
        if image.mode != 'RGB':
            image = image.convert('RGB')

        # Perform OCR on the image
        result = ocr.ocr(image, cls=True)

        # Extract text from the OCR result
        text = ""
        for line in result[0]:
            text += line[1] + " "
        return text
    except Exception as e:
        print(f"Error processing image: {e}")
        return ""


def ask_gemini_ai_with_genai(text: str):
    """Generate a summary and flashcards using Google Generative AI."""
    genai.configure(api_key="AIzaSyAe9JnvjQNtr95ZeLEXLS9Ay9yzrA7vidM")
    model = genai.GenerativeModel('gemini-pro')

    # Generate content
    prompt = (
        f"Summarize the following content in paragraph form :\n\n{text}"
    )
    prompt2 = (
        f"Generate flashcard with question and answer: in paragraph form\n\n{text}"
    )
    response = model.generate_content(prompt)
    response2 = model.generate_content(prompt2)

    # Parse the generated content
    generated_text = response.text
    print(generated_text)
    summary = [response.text]
    flashcards = [response2.text]

    print("Summary:", summary)
    print("FLASHCARD:", flashcards)
    return summary, flashcards


def ask_gemini_for_answer(question: str):
    """Query Gemini AI for an answer to a specific question."""
    genai.configure(api_key="AIzaSyAe9JnvjQNtr95ZeLEXLS9Ay9yzrA7vidM")
    model = genai.GenerativeModel('gemini-pro')

    # Generate answer to the question
    prompt = f"Answer the following question:\n\n{question}"

    response = model.generate_content(prompt)

    # Return the response text
    return response.text

@app.route('/')
def index():
    return render_template('index.html')


@app.route('/ask_question', methods=['POST'])
def ask_question():
    """Endpoint for asking a custom question to Gemini AI."""
    question = request.form.get('question')

    if not question:
        return jsonify({"error": "No question provided"}), 400

    # Get the answer from Gemini AI
    answer = ask_gemini_for_answer(question)

    return jsonify({
        "question": question,
        "answer": answer
    })
@app.route('/generate_flashcards', methods=['POST'])
def generate_flashcards():
    if 'file' not in request.files:
        return jsonify({"error": "No file uploaded"}), 400

    file = request.files['file']
    if file.filename == '':
        return jsonify({"error": "No selected file"}), 400

    if file and file.filename.endswith('.pdf'):
        text = extract_text_from_pdf(file)
    elif file and file.filename.endswith('.docx'):
        text = extract_text_from_docx(file)
    elif file and file.filename.endswith('.pptx'):
        text = extract_text_from_pptx(file)
    elif file and file.filename.lower().endswith(('.png', '.jpg', '.jpeg')):
        text = extract_text_from_image(file)
    else:
        return jsonify({"error": "Unsupported file format"}), 400

    summary, flashcards = ask_gemini_ai_with_genai(text)

    # Save to Python lists or database
    tracked["summary"] = tracked.get("summary", []) + summary
    tracked["flashcards"] = tracked.get("flashcards", []) + flashcards
    save_tracked()

    return jsonify({
        "summary_points": summary,
        "flashcards": flashcards,
        "store_in_local": True  # Suggests storing data in local storage for subsequent use
    })


if __name__ == "__main__":
    app.run(debug=True)
